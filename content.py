from pptx import Presentation

# Load the PPT file
presentation = Presentation('scores.xlsx')

# Extract names and scores from slides
data = []
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:  # Add condition to filter names and scores (depends on the PPT format)
                parts = text.split()  # Assuming name and score are separated by whitespace
                if len(parts) >= 2:  # Name and score detected
                    name = ' '.join(parts[:-1])  # All but last is name
                    score = parts[-1]           # Last part is score
                    data.append((name, score))

# Generate HTML
html_content = '<ol>\n'
for item in data:
    name, score = item
    html_content += f'  <li>\n    <p class="name">{name}</p>\n    <p class="score">{score}</p>\n  </li>\n'
html_content += '</ol>'

# Save HTML to a file
with open('output.html', 'w') as f:
    f.write(html_content)

print("HTML generated successfully!")
